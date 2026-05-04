from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.util import Inches, Pt
import copy

# ─── THEME COLORS ────────────────────────────────────────────────────────────
C_DARK_BG    = RGBColor(0x0D, 0x1B, 0x2A)   # Deep navy background
C_CARD_BG    = RGBColor(0x11, 0x29, 0x40)   # Slightly lighter navy for cards
C_ACCENT     = RGBColor(0x4F, 0x46, 0xE5)   # Indigo accent (matches app)
C_ACCENT2    = RGBColor(0x06, 0xB6, 0xD4)   # Cyan accent
C_WHITE      = RGBColor(0xFF, 0xFF, 0xFF)
C_LIGHT_GRAY = RGBColor(0x94, 0xA3, 0xB8)
C_GOLD       = RGBColor(0xF5, 0x9E, 0x0B)

W = Inches(13.33)  # Widescreen width
H = Inches(7.5)    # Widescreen height

prs = Presentation()
prs.slide_width  = W
prs.slide_height = H

blank_layout = prs.slide_layouts[6]  # Completely blank

# ─── HELPERS ─────────────────────────────────────────────────────────────────

def add_rect(slide, x, y, w, h, color, alpha=None):
    shape = slide.shapes.add_shape(1, x, y, w, h)
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    return shape

def add_text(slide, text, x, y, w, h,
             size=18, bold=False, color=C_WHITE,
             align=PP_ALIGN.LEFT, italic=False, wrap=True):
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = wrap
    p = tf.paragraphs[0]
    p.alignment = align
    run = p.add_run()
    run.text = text
    run.font.size = Pt(size)
    run.font.bold = bold
    run.font.italic = italic
    run.font.color.rgb = color
    run.font.name = "Calibri"
    return txBox

def add_slide(title_text, subtitle_text=None):
    slide = prs.slides.add_slide(blank_layout)

    # Full background
    add_rect(slide, 0, 0, W, H, C_DARK_BG)

    # Left accent bar
    add_rect(slide, 0, 0, Inches(0.12), H, C_ACCENT)

    # Top accent strip
    add_rect(slide, Inches(0.12), 0, W - Inches(0.12), Inches(0.08), C_ACCENT2)

    # Bottom bar
    add_rect(slide, 0, H - Inches(0.45), W, Inches(0.45), C_CARD_BG)
    add_text(slide, "NextGenISP  |  ISP Management Platform  |  Confidential",
             Inches(0.3), H - Inches(0.38), W - Inches(0.6), Inches(0.35),
             size=9, color=C_LIGHT_GRAY, align=PP_ALIGN.LEFT)
    add_text(slide, "2026",
             0, H - Inches(0.38), W - Inches(0.3), Inches(0.35),
             size=9, color=C_LIGHT_GRAY, align=PP_ALIGN.RIGHT)

    # Slide title pill background
    pill = add_rect(slide, Inches(0.3), Inches(0.18), Inches(9), Inches(0.72), C_ACCENT)
    # Title text
    add_text(slide, title_text, Inches(0.45), Inches(0.2), Inches(9), Inches(0.68),
             size=26, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)

    if subtitle_text:
        add_text(slide, subtitle_text, Inches(0.45), Inches(0.95), W - Inches(0.9), Inches(0.5),
                 size=14, bold=False, color=C_ACCENT2, align=PP_ALIGN.LEFT, italic=True)

    return slide

def add_body_text(slide, lines, x=Inches(0.4), y=Inches(1.6), w=None, h=Inches(5.3),
                  size=14, bullet=True, color=C_WHITE, spacing=1.2):
    if w is None:
        w = W - Inches(0.8)
    txBox = slide.shapes.add_textbox(x, y, w, h)
    tf = txBox.text_frame
    tf.word_wrap = True
    first = True
    for line in lines:
        if first:
            p = tf.paragraphs[0]
            first = False
        else:
            p = tf.add_paragraph()
        p.alignment = PP_ALIGN.LEFT
        p.space_before = Pt(4)
        p.space_after  = Pt(2)
        run = p.add_run()
        if isinstance(line, tuple):
            txt, sz, bd, clr = line
            run.text = txt
            run.font.size = Pt(sz)
            run.font.bold = bd
            run.font.color.rgb = clr
        else:
            run.text = ("• " if bullet else "") + line
            run.font.size = Pt(size)
            run.font.color.rgb = color
        run.font.name = "Calibri"

def two_column(slide, left_lines, right_lines,
               lx=Inches(0.4), ly=Inches(1.55), lw=Inches(6),
               rx=Inches(6.9), ry=Inches(1.55), rw=Inches(6),
               h=Inches(5.4), size=13):
    add_body_text(slide, left_lines,  x=lx, y=ly, w=lw, h=h, size=size)
    add_body_text(slide, right_lines, x=rx, y=ry, w=rw, h=h, size=size)

def section_header(slide, text, y=Inches(1.45)):
    add_rect(slide, Inches(0.4), y, Inches(0.06), Inches(0.32), C_GOLD)
    add_text(slide, text, Inches(0.55), y, Inches(10), Inches(0.35),
             size=15, bold=True, color=C_GOLD, align=PP_ALIGN.LEFT)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 1 — TITLE
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, H, C_DARK_BG)
add_rect(slide, 0, 0, Inches(0.12), H, C_ACCENT)
add_rect(slide, Inches(0.12), 0, W, Inches(0.08), C_ACCENT2)
add_rect(slide, 0, H - Inches(0.45), W, Inches(0.45), C_CARD_BG)

# Big glowing center block
add_rect(slide, Inches(1.5), Inches(1.6), Inches(10.3), Inches(4), C_CARD_BG)
add_rect(slide, Inches(1.5), Inches(1.6), Inches(0.18), Inches(4), C_ACCENT)

add_text(slide, "NextGenISP", Inches(1.9), Inches(1.8), Inches(9.5), Inches(1.6),
         size=68, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
add_text(slide, "An Integrated ERP & CRM Platform for Internet Service Providers",
         Inches(1.9), Inches(3.3), Inches(9.5), Inches(0.8),
         size=20, bold=False, color=C_ACCENT2, align=PP_ALIGN.LEFT, italic=True)
add_text(slide, "Built with  Django REST Framework  ·  React (Vite)  ·  Leaflet.js  ·  Recharts",
         Inches(1.9), Inches(4.05), Inches(9.5), Inches(0.6),
         size=14, color=C_LIGHT_GRAY, align=PP_ALIGN.LEFT)
add_text(slide, "Academic Project Presentation  |  2026",
         Inches(1.9), Inches(4.7), Inches(9.5), Inches(0.5),
         size=12, color=C_LIGHT_GRAY, align=PP_ALIGN.LEFT)
add_text(slide,
         "NextGenISP  |  ISP Management Platform  |  Confidential",
         Inches(0.3), H - Inches(0.38), W - Inches(0.6), Inches(0.35),
         size=9, color=C_LIGHT_GRAY, align=PP_ALIGN.LEFT)
add_text(slide, "2026", 0, H - Inches(0.38), W - Inches(0.3), Inches(0.35),
         size=9, color=C_LIGHT_GRAY, align=PP_ALIGN.RIGHT)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 2 — INTRODUCTION
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide("Introduction", "What is NextGenISP?")

add_body_text(slide, [
    ("What is NextGenISP?", 16, True, C_ACCENT2),
    "NextGenISP is a full-stack ISP Management Information System (MIS) built to digitize and automate all operations of an Internet Service Provider.",
    "",
    ("ERP — Enterprise Resource Planning", 15, True, C_GOLD),
    "Automates all internal business processes into a single platform.",
    "  → Resource Mgmt: Hardware inventory, MAC address binding, device tracking",
    "  → Workforce Planning: Dual-role staff (Field + Technical) with Workload Monitor",
    "  → Financial Ops: Invoice generation, billing cycles, payment gateway",
    "  → Infrastructure: GeoJSON zone management, NOC live map, outage control",
    "",
    ("CRM — Customer Relationship Management", 15, True, C_GOLD),
    "Manages the complete subscriber lifecycle from lead to loyal customer.",
    "  → Lead Management: LEAD → VERIFIED → ACTIVE onboarding pipeline",
    "  → Self-Service Portal: Speed test, diagnostics, bills, AI chatbot",
    "  → Support & Ticketing: Physical & Logical issue routing to right staff",
    "  → Broadcast Engine: Zone-targeted email campaigns for outage alerts",
], size=13)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 3 — ABSTRACT
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide("Abstract", "Project at a Glance")

add_body_text(slide, [
    "Traditional ISPs manage operations through spreadsheets, WhatsApp, and manual processes — creating data silos, billing errors, and poor customer experience.",
    "",
    ("NextGenISP solves this by providing:", 15, True, C_ACCENT2),
    "  ✦  Public Portal — Self-service plan selection & new connection registration",
    "  ✦  Admin Command Center — Full business management in one dashboard",
    "  ✦  Dual-Staff Dashboards — Physical install (Field) + Logical config (Technical)",
    "  ✦  Customer Self-Service — Speed gauge, diagnostic terminal, billing, AI chatbot",
    "  ✦  Live NOC War Room — GeoJSON zone map with real-time customer pins",
    "  ✦  Targeted Broadcasts — Zone/role-based email alerts with customer name preview",
    "",
    ("Key Innovations:", 15, True, C_GOLD),
    "  →  Animated SVG Speedometer simulating plan-based real speeds",
    "  →  Hacker-style Diagnostic Terminal with auto ticket raising",
    "  →  Simulated Razorpay + UPI QR payment gateway",
    "  →  AI Chatbot using JavaScript decision-tree for instant FAQ resolution",
], size=13)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 4 — PROBLEM STATEMENT
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide("Problem Statement", "The Pain Points of Traditional ISPs")

rows = [
    ("Customer data in Excel / WhatsApp",      "No traceability, data loss risk"),
    ("No formal ticket/complaint system",       "Complaints ignored or delayed"),
    ("Field staff managed via phone calls",     "No tracking, wasted site visits"),
    ("Manual invoice generation",               "Billing errors, late payments"),
    ("No geographic network visibility",        "Cannot identify area-wise outages"),
    ("Hardware not tracked post-delivery",      "No MAC binding, security risk"),
    ("Customers call for basic queries",        "High support load, slow response"),
    ("No centralized broadcast capability",     "Cannot alert customers during outages"),
]

# Table header
add_rect(slide, Inches(0.35), Inches(1.5), Inches(6.2), Inches(0.38), C_ACCENT)
add_rect(slide, Inches(6.6), Inches(1.5), Inches(6.4), Inches(0.38), C_ACCENT2)
add_text(slide, "Problem", Inches(0.45), Inches(1.52), Inches(6), Inches(0.36),
         size=13, bold=True, color=C_WHITE)
add_text(slide, "Impact", Inches(6.7), Inches(1.52), Inches(6.2), Inches(0.36),
         size=13, bold=True, color=C_DARK_BG)

row_h = Inches(0.52)
for i, (prob, impact) in enumerate(rows):
    y = Inches(1.88) + i * row_h
    bg = C_CARD_BG if i % 2 == 0 else C_DARK_BG
    add_rect(slide, Inches(0.35), y, Inches(6.2), row_h - Inches(0.02), bg)
    add_rect(slide, Inches(6.6), y, Inches(6.4), row_h - Inches(0.02), bg)
    add_text(slide, f"✗  {prob}", Inches(0.45), y + Inches(0.05), Inches(6), row_h,
             size=12, color=RGBColor(0xFF, 0x99, 0x99))
    add_text(slide, impact, Inches(6.7), y + Inches(0.05), Inches(6.2), row_h,
             size=12, color=C_LIGHT_GRAY)

add_rect(slide, Inches(0.35), Inches(6.2), Inches(12.65), Inches(0.05), C_ACCENT)
add_text(slide,
         "There is a need for a unified digital system managing all ISP operations — from customer onboarding to network monitoring.",
         Inches(0.4), Inches(6.25), Inches(12.6), Inches(0.5),
         size=13, bold=True, color=C_GOLD, align=PP_ALIGN.CENTER)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 5 — EXISTING VS PROPOSED
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide("Existing System  vs  Proposed System")

features = [
    ("Customer Registration", "Paper forms / walk-in", "Online portal + ID proof upload"),
    ("Plan Selection",         "Phone call with agent", "Self-service catalog with filters"),
    ("Installation Tracking",  "WhatsApp messages",     "Two-phase digital task board"),
    ("Billing",                "Manual invoice printing","Auto-invoices + payment gateway"),
    ("Support Tickets",        "Verbal, no tracking",   "Typed tickets with auto-routing"),
    ("Network Monitoring",     "Engineers go on-site",  "Live Leaflet.js NOC map"),
    ("Staff Workload",         "No visibility",         "Workload charts per staff member"),
    ("Mass Communication",     "Manual SMS blast",      "Targeted zone-aware broadcast"),
    ("Hardware Tracking",      "No records",            "MAC binding + inventory system"),
    ("Customer Self-Service",  "None",                  "Speed test, diagnostics, chatbot"),
]

hdrs = [Inches(0.35), Inches(4.45), Inches(8.9)]
ws   = [Inches(4),    Inches(4.35),  Inches(4.35)]
hdr_labels = ["Feature", "Existing (Manual)", "Proposed (NextGenISP)"]
hdr_colors = [C_CARD_BG, RGBColor(0x7F, 0x1D, 0x1D), C_ACCENT]

for j, (lbl, hx, hw) in enumerate(zip(hdr_labels, hdrs, ws)):
    add_rect(slide, hx, Inches(1.45), hw, Inches(0.38), hdr_colors[j])
    add_text(slide, lbl, hx + Inches(0.07), Inches(1.47), hw, Inches(0.36),
             size=12, bold=True, color=C_WHITE)

row_h = Inches(0.48)
for i, (feat, old, new) in enumerate(features):
    y = Inches(1.83) + i * row_h
    bg = C_CARD_BG if i % 2 == 0 else C_DARK_BG
    for hx, hw in zip(hdrs, ws):
        add_rect(slide, hx, y, hw, row_h - Inches(0.02), bg)
    vals = [feat, old, new]
    clrs = [C_WHITE, RGBColor(0xFF, 0x80, 0x80), RGBColor(0x6E, 0xFF, 0xB0)]
    for hx, hw, val, clr in zip(hdrs, ws, vals, clrs):
        add_text(slide, val, hx + Inches(0.07), y + Inches(0.05), hw - Inches(0.1), row_h,
                 size=11.5, color=clr)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 6 — DFD LEVELS
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide("Data Flow Diagram (DFD)", "Levels of System Data Flow")

section_header(slide, "Level 0 — Context Diagram (The Big Picture)", y=Inches(1.45))
add_body_text(slide, [
    "The entire system is a single process: \"NextGenISP Platform\".",
    "External Entities:  Visitor  ·  Customer  ·  Admin  ·  Field Staff  ·  Technical Staff  ·  Email System",
], x=Inches(0.5), y=Inches(1.82), h=Inches(1.1), size=13, bullet=False)

section_header(slide, "Level 1 — Main Processes (6 Core Functions)", y=Inches(2.95))
add_body_text(slide, [
    "  [1.0]  User & Auth Management   → JWT Token, Role-locked Dashboard",
    "  [2.0]  Plan & Hardware Mgmt     → Public plan listing, Subscriptions",
    "  [3.0]  Installation Management  → Field task board, MAC binding",
    "  [4.0]  Billing & Payments       → Invoice generation, Payment records",
    "  [5.0]  Support & Ticketing      → Ticket creation, Staff routing, Resolution",
    "  [6.0]  NOC / Network Ops        → Maintenance alerts, Live map visualization",
], x=Inches(0.5), y=Inches(3.3), h=Inches(2.0), size=13, bullet=False)

section_header(slide, "Level 2 — Installation Sub-Process (Drill-Down)", y=Inches(5.35))
add_body_text(slide, [
    "  [3.1] Admin creates task  →  [3.2] Field Staff does physical install  →  [3.3] Tech Staff does logical config  →  [3.4] Customer ACTIVE",
], x=Inches(0.5), y=Inches(5.65), h=Inches(0.8), size=12, bullet=False, color=C_ACCENT2)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 7 — ER DIAGRAM DESCRIPTION
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide("ER Diagram", "Entity-Relationship Description")

add_body_text(slide, [
    ("Core Entities & Relationships", 15, True, C_ACCENT2),
    "",
    ("Area  (Geographic anchor)", 13, True, C_GOLD),
    "  → User (Customer) belongs to one Area  [Many-to-One FK]",
    "  → Plan is available in multiple Areas  [Many-to-Many]",
    "",
    ("User  (Most connected entity)", 13, True, C_GOLD),
    "  → Holds a Subscription → Plan  [FK chain]",
    "  → Subscription generates monthly Invoices  [One-to-Many]",
    "  → Raises Tickets assigned to Staff  [Two FKs on User table]",
    "  → Linked to InstallationTask as Customer, Field Staff, OR Tech Staff  [×3 FKs]",
    "",
    ("InstallationTask  (Operations hub)", 13, True, C_GOLD),
    "  → References: Customer, Field Staff, Tech Staff, Plan, Hardware",
    "  → Stores router MAC address, own-device model, one-time charges",
    "",
    ("Supporting entities: Ticket → ChatMessage  |  InventoryItem → AssetAssignment",
     12, False, C_LIGHT_GRAY),
], size=13, bullet=False)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 8 — DATABASE SCHEMA  (split into 2 slides)
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide("Database Schema — Part 1", "13 Tables across 8 Groups")

tables_left = [
    ("area", ["id PK", "name, code (UNIQUE)", "coordinates (GeoJSON)", "is_under_maintenance", "maintenance_message"]),
    ("user  (extends AbstractUser)", ["id PK", "username, email, password", "role  ENUM: ADMIN|TECH|FIELD|CUSTOMER", "status  ENUM: LEAD→ACTIVE→SUSPENDED", "area_id FK, phone_number, address", "latitude, longitude, id_proof"]),
    ("plan", ["id PK", "name, speed_mbps, price", "plan_type ENUM: FIBER|WIRELESS", "data_limit_gb (FUP cap)", "areas  M2M→ area", "recommended_hardware  M2M→ hardware"]),
]
tables_right = [
    ("subscription", ["user_id FK", "plan_id FK", "status ENUM: ACTIVE|EXPIRED|SUSPENDED", "start_date, end_date, billing_cycle"]),
    ("invoice", ["subscription_id FK", "amount, due_date", "status ENUM: PENDING|PAID|OVERDUE"]),
    ("payment", ["invoice_id FK", "transaction_id (UNIQUE)", "amount, method (CASH|ONLINE)"]),
    ("hardware", ["id PK", "name, price", "specifications (JSON)", "image, is_active"]),
]

def render_table_card(slide, table_name, fields, x, y, w=Inches(5.9)):
    h_hdr = Inches(0.3)
    h_row = Inches(0.22)
    total_h = h_hdr + len(fields) * h_row + Inches(0.08)
    add_rect(slide, x, y, w, h_hdr, C_ACCENT)
    add_text(slide, f"  {table_name}", x, y, w, h_hdr + Inches(0.02),
             size=11, bold=True, color=C_WHITE)
    add_rect(slide, x, y + h_hdr, w, total_h - h_hdr, C_CARD_BG)
    for i, f in enumerate(fields):
        add_text(slide, f"    {f}", x, y + h_hdr + i * h_row, w, h_row + Inches(0.04),
                 size=10, color=C_LIGHT_GRAY)
    return total_h

y = Inches(1.52)
gap = Inches(0.15)
for tname, fields in tables_left:
    h = render_table_card(slide, tname, fields, Inches(0.35), y)
    y += h + gap

y = Inches(1.52)
for tname, fields in tables_right:
    h = render_table_card(slide, tname, fields, Inches(6.75), y)
    y += h + gap

# ─── DB Schema Part 2 ────────────────────────────────────────────────────────
slide = add_slide("Database Schema — Part 2", "Remaining Tables & Relationships")

tables_left2 = [
    ("ticket", ["customer_id FK → user", "assigned_to_id FK → user", "ticket_type ENUM: LOGICAL|PHYSICAL|INSTALLATION", "status ENUM: OPEN→IN_PROGRESS→RESOLVED→CLOSED", "subject, description"]),
    ("chatmessage", ["ticket_id FK", "sender_id FK → user", "message, created_at"]),
    ("enquiry", ["name, phone, email", "status ENUM: OPEN→CONTACTED→CONVERTED", "assigned_to_id FK → user"]),
]
tables_right2 = [
    ("installationtask", ["customer_id FK → user", "assigned_staff_id FK → user  (Field)", "assigned_technical_staff_id FK (Tech)", "plan_id FK, hardware_id FK", "own_router_mac, own_router_model", "status ENUM: PENDING→PHYSICAL_COMPLETED→CLOSED", "one_time_charge, router_mac, notes"]),
    ("inventoryitem", ["name, sku (UNIQUE)", "quantity, unit_price", "category: ROUTER|CABLE|CONNECTOR"]),
    ("assetassignment", ["item_id FK", "user_id FK", "serial_number, is_returned"]),
]

y = Inches(1.52)
for tname, fields in tables_left2:
    h = render_table_card(slide, tname, fields, Inches(0.35), y)
    y += h + gap

y = Inches(1.52)
for tname, fields in tables_right2:
    h = render_table_card(slide, tname, fields, Inches(6.75), y)
    y += h + gap

add_rect(slide, Inches(0.35), Inches(6.2), Inches(12.65), Inches(0.6), C_CARD_BG)
add_text(slide,
         "area ← user → subscription → invoice → payment  |  user → ticket → chatmessage  |  user → installationtask  |  inventoryitem → assetassignment",
         Inches(0.45), Inches(6.22), Inches(12.5), Inches(0.55),
         size=10, color=C_ACCENT2, bold=False, italic=True)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 9 — MODULE DESCRIPTIONS
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide("Module Descriptions", "5 Core Modules of the System")

modules = [
    ("01", "Public Portal",
     "Visitors browse plans, hardware, and register. Account created as LEAD for admin verification.",
     "Home · Plan List · Router Detail · Register · Login"),
    ("02", "Admin Dashboard",
     "Verify users, manage plans/zones/hardware, view workload charts, send targeted broadcasts.",
     "AdminDash · ManageUsers · ManagePlans · ManageAreas · Broadcast · WorkloadMonitor"),
    ("03", "War Room (NOC)",
     "Live Leaflet.js map with GeoJSON zone polygons. Color-coded customer pins. Toggle maintenance mode per zone.",
     "NetworkMap · Outage Tracker"),
    ("04", "Staff Dashboards (Two-Phase)",
     "Field Staff: physical cabling, MAC capture. Technical Staff: logical config (PPPoE/DHCP/VLAN), KB articles.",
     "FieldDash · FieldRepairs · TechDash · TechRepairs · KnowledgeBase"),
    ("05", "Customer Self-Service Portal",
     "Plan status, speed gauge, diagnostic terminal, bill history, UPI/Razorpay payment, tickets, AI chatbot.",
     "CustDash · BillHistory · Support · Chatbot"),
]

y = Inches(1.5)
for num, mod, desc, pages in modules:
    add_rect(slide, Inches(0.35), y, Inches(0.55), Inches(0.9), C_ACCENT)
    add_text(slide, num, Inches(0.35), y, Inches(0.55), Inches(0.9),
             size=14, bold=True, color=C_WHITE, align=PP_ALIGN.CENTER)
    add_rect(slide, Inches(0.95), y, Inches(12.05), Inches(0.9), C_CARD_BG)
    add_text(slide, mod, Inches(1.08), y + Inches(0.04), Inches(6), Inches(0.36),
             size=14, bold=True, color=C_GOLD)
    add_text(slide, desc, Inches(1.08), y + Inches(0.38), Inches(7.5), Inches(0.35),
             size=11, color=C_WHITE)
    add_text(slide, f"Pages: {pages}", Inches(1.08), y + Inches(0.65), Inches(11.7), Inches(0.26),
             size=9.5, color=C_LIGHT_GRAY, italic=True)
    y += Inches(1.0)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 10 — ADVANTAGES
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide("Advantages of the Project")

advantages = [
    ("End-to-End Digitization",      "No paperwork — entire ISP workflow is digital from lead to active subscriber"),
    ("Real-Time Visibility",         "Admin sees live NOC map, workload charts, outages, and revenue at a glance"),
    ("Dual-Staff Workflow",          "Physical and logical roles are separated, eliminating assignment confusion"),
    ("Self-Service Empowerment",     "Customers resolve issues themselves — speed test, diagnostics, chatbot"),
    ("Targeted Communications",      "Zone-level broadcast ensures only affected users receive alerts"),
    ("Hardware Accountability",      "MAC binding and inventory tracking prevent unauthorized device usage"),
    ("Scalable Role Architecture",   "New roles can be added via DB config — zero code changes needed"),
    ("Geographic Intelligence",      "GeoJSON polygon zones tie plans, customers, and alerts to geography"),
    ("Automated Billing",            "Invoice generation, payment tracking, defaulter detection — fully automated"),
    ("Secure & Authenticated",       "JWT-based auth + role-locked API endpoints at every backend view"),
]

cols = [advantages[:5], advantages[5:]]
xs = [Inches(0.35), Inches(6.75)]
for col, x in zip(cols, xs):
    y = Inches(1.52)
    for title, desc in col:
        add_rect(slide, x, y, Inches(6.2), Inches(0.9), C_CARD_BG)
        add_rect(slide, x, y, Inches(0.08), Inches(0.9), C_ACCENT2)
        add_text(slide, f"✓  {title}", x + Inches(0.18), y + Inches(0.05), Inches(5.9), Inches(0.35),
                 size=13, bold=True, color=C_ACCENT2)
        add_text(slide, desc, x + Inches(0.18), y + Inches(0.42), Inches(5.9), Inches(0.42),
                 size=11, color=C_LIGHT_GRAY)
        y += Inches(1.0)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 11 — FUTURE SCOPE
# ═══════════════════════════════════════════════════════════════════════════════
slide = add_slide("Future Scope", "What Comes Next")

section_header(slide, "Standard Enhancements", y=Inches(1.45))
add_body_text(slide, [
    "  ✦  Real Razorpay Integration — Webhook-based payment confirmation for production",
    "  ✦  Progressive Web App (PWA) — Field Staff access tasks offline on mobile",
    "  ✦  WhatsApp / SMS Notifications — Instant payment reminders via Twilio",
    "  ✦  Usage-Based FUP Enforcement — Auto-suspend on data cap breach",
], x=Inches(0.5), y=Inches(1.82), h=Inches(1.5), size=13, bullet=False)

section_header(slide, "AI / Machine Learning Enhancements", y=Inches(3.38))
add_body_text(slide, [
    "  🤖  Predictive Maintenance (ML) — Forecast fiber degradation by zone before customers notice",
    "  🤖  Auto Ticket Triage (LLM) — Classify complaints as LOGICAL/PHYSICAL and assign right staff",
    "  🤖  Dynamic Bandwidth QoS (AI) — Prioritize premium customers during peak hours automatically",
    "  🤖  Churn Prediction — Identify customers likely to cancel and trigger retention campaigns",
], x=Inches(0.5), y=Inches(3.75), h=Inches(1.65), size=13, bullet=False)

section_header(slide, "Infrastructure Scaling", y=Inches(5.42))
add_body_text(slide, [
    "  ⚙  PostgreSQL Migration  ·  Docker Containerization  ·  WebSocket Real-Time Ticket Chat",
], x=Inches(0.5), y=Inches(5.76), h=Inches(0.6), size=13, bullet=False, color=C_LIGHT_GRAY)

# ═══════════════════════════════════════════════════════════════════════════════
# SLIDE 12 — THANK YOU
# ═══════════════════════════════════════════════════════════════════════════════
slide = prs.slides.add_slide(blank_layout)
add_rect(slide, 0, 0, W, H, C_DARK_BG)
add_rect(slide, 0, 0, Inches(0.12), H, C_ACCENT)
add_rect(slide, Inches(0.12), 0, W, Inches(0.08), C_ACCENT2)
add_rect(slide, 0, H - Inches(0.45), W, Inches(0.45), C_CARD_BG)
add_rect(slide, Inches(1.5), Inches(1.8), Inches(10.3), Inches(3.5), C_CARD_BG)
add_rect(slide, Inches(1.5), Inches(1.8), Inches(0.18), Inches(3.5), C_ACCENT2)
add_text(slide, "Thank You", Inches(1.9), Inches(2.0), Inches(9.5), Inches(1.4),
         size=60, bold=True, color=C_WHITE, align=PP_ALIGN.LEFT)
add_text(slide, "Questions & Discussion", Inches(1.9), Inches(3.35), Inches(9.5), Inches(0.6),
         size=22, color=C_ACCENT2, italic=True, align=PP_ALIGN.LEFT)
add_text(slide, "NextGenISP  ·  Built with Django + React  ·  2026",
         Inches(1.9), Inches(4.0), Inches(9.5), Inches(0.5),
         size=13, color=C_LIGHT_GRAY, align=PP_ALIGN.LEFT)
add_text(slide, "NextGenISP  |  ISP Management Platform  |  Confidential",
         Inches(0.3), H - Inches(0.38), W - Inches(0.6), Inches(0.35),
         size=9, color=C_LIGHT_GRAY, align=PP_ALIGN.LEFT)

# ─── SAVE ────────────────────────────────────────────────────────────────────
output_path = r"E:\Backup\Main  projecct\NextGenISP_Presentation.pptx"
prs.save(output_path)
print(f"\n✅  Presentation saved to:\n    {output_path}\n")
print(f"    Slides: {len(prs.slides)}")
print("    All done!")
